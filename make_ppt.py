"""Build the presentation deck for the pith-position exercise (16:9, python-pptx)."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent
OUT = ROOT / "outputs"

DARK = RGBColor(0x1F, 0x2D, 0x3D)
ACCENT = RGBColor(0x2E, 0x6B, 0xB0)
GRAY = RGBColor(0x5A, 0x64, 0x72)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title=None, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.85))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = DARK
        if subtitle:
            p2 = tb.text_frame.add_paragraph()
            r = p2.add_run(); r.text = subtitle
            r.font.size = Pt(14); r.font.color.rgb = GRAY
    return s


def bullets(slide, items, left=0.6, top=1.3, width=12.1, height=5.8, size=17):
    height = min(height, 7.4 - top)
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size if level == 0 else size - 2)
        r.font.color.rgb = DARK if level == 0 else GRAY
    return tb


def picture(slide, path, left, top, width=None, height=None):
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


# ---------------- 1. Title ----------------
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(2.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Predicting the Pith Position from a Board-Face Image"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK
p2 = tb.text_frame.add_paragraph()
r = p2.add_run(); r.text = "A compact CNN for 15-point pith x-regression  ·  Test MAE 3.82 mm"
r.font.size = Pt(20); r.font.color.rgb = ACCENT
p3 = tb.text_frame.add_paragraph()
r = p3.add_run(); r.text = "Candidate exercise — Doctoral Studentship in Building Technology (AI-Based Timber Quality Assessment)"
r.font.size = Pt(14); r.font.color.rgb = GRAY

# ---------------- 2. Introduction ----------------
s = add_slide("Introduction: why locate the pith?")
bullets(s, [
    ("The pith is the centre of the tree stem; growth rings form around it.", 0),
    ("Its position relative to a sawn board tells us where the board lay in the log — useful for wood characterisation, sawing decisions and strength grading.", 0),
    ("Task in this exercise (simplified from the full 4-face, x+y problem):", 0),
    ("Input: one RGB image (512×512) of one wide board face, 145 mm × 145 mm.", 1),
    ("Output: pith x-position (mm) at 15 equally spaced z positions, bottom → top.", 1),
    ("The pith may lie outside the board — predictions must not be clamped to [0, 145] mm.", 1),
    ("Metric: mean absolute error (mm) over all 15 positions × 1000 test boards.", 0),
    ("Data: 4000 labelled training boards, 1000 held-out test boards (split kept fixed; test never used for training).", 0),
], top=1.35)

# ---------------- 3. Data exploration ----------------
s = add_slide("First look at the data", "Ground-truth pith path (red) drawn over training images")
picture(s, OUT / "explore_samples.png", 0.35, 1.35, width=12.6)
bullets(s, [
    ("Pith inside the board → rings form a clear symmetry axis with high curvature around it (left two).", 0),
    ("Pith far outside (3rd image, x ≈ 215 mm) → rings become almost straight vertical stripes; position must be inferred from weak curvature and ring-spacing gradients.", 0),
    ("Knots, colour variation and blur act as realistic nuisance factors.", 0),
], top=4.85, size=15)

# ---------------- 4. Label statistics ----------------
s = add_slide("What the labels say", "Statistics over the 4000 training boards")
picture(s, OUT / "label_stats.png", 0.35, 1.35, width=12.6)
bullets(s, [
    ("x spans −60 … +221 mm; 4.6% of boards have the pith outside the face for at least one z — a rare but important sub-population.", 0),
    ("The pith path along z is smooth: median within-board wander is only 10.4 mm → z-neighbouring outputs are strongly correlated.", 0),
    ("Naive baseline (always predict the training mean per position): 25.6 mm test MAE — the number any model must beat.", 0),
], top=5.0, size=15)

# ---------------- 5. Formulation & choices ----------------
s = add_slide("Problem formulation and key choices")
bullets(s, [
    ("Formulation: direct regression of 15 continuous values from one image (no clamping, so out-of-board positions are handled naturally).", 0),
    ("Alternative considered — 1D heatmap + soft-argmax per z position: more mechanics for little expected gain at this scale; kept as future work.", 1),
    ("Input resolution 224×224 (from 512): ring curvature and spacing cues survive Lanczos downscaling; ~5× faster than 512² on the available hardware.", 0),
    ("Compact ~1.0 M-parameter CNN trained from scratch, rather than a pretrained ImageNet backbone:", 0),
    ("single synthetic domain + 4000 samples → pretraining buys little; measured 4× faster per epoch than ResNet-18 on the CPU-only machine (2 vs 8 min).", 1),
    ("Targets normalised t = (x − 72.5)/145; SmoothL1 (Huber) loss — robust to the few extreme out-of-board targets.", 0),
    ("Augmentation uses the two exact physical symmetries: horizontal flip (x → 145 − x) and vertical flip (reverse the 15 targets). No unrealistic augmentations.", 0),
    ("All model selection on a 3600/400 train/validation split — the 1000 test boards were touched exactly once, at the end.", 0),
], top=1.3, size=16)

# ---------------- 6. Architecture ----------------
s = add_slide("Architecture: pool over width, keep z structure")
picture(s, OUT / "architecture.png", 0.35, 1.5, width=12.6)
bullets(s, [
    ("The output has spatial structure: target i depends mostly on rings near z-band i.", 0),
    ("So the head pools the feature map over the width only → 15 vertical bands, each mapped to one x-value by a shared 1×1 convolution (same weights per band).", 0),
    ("Deep layers still see the whole image (receptive field), so global context is not lost — the head just aligns parameters with the geometry of the task.", 0),
], top=4.9, size=15)

# ---------------- 7. Training ----------------
s = add_slide("Training setup and convergence")
picture(s, OUT / "training_curve_final.png", 6.7, 1.4, width=6.2)
bullets(s, [
    ("AdamW, lr 3·10⁻⁴, weight decay 10⁻⁴, cosine schedule.", 0),
    ("Batch 32, 35 epochs, SmoothL1 loss on normalised targets.", 0),
    ("Best checkpoint chosen by validation MAE (4.25 mm at epoch 30).", 0),
    ("Whole pipeline runs on a 4-core laptop CPU:", 0),
    ("images pre-decoded once to a uint8 cache → no PNG decoding during training,", 1),
    ("~2 min/epoch → ~65 min total training time.", 1),
    ("Deterministic seeds for split and initialisation → reproducible.", 0),
], left=0.6, top=1.5, width=5.9, size=16)

# ---------------- 8. Ablation ----------------
s = add_slide("Does the structured head help? A controlled ablation")
picture(s, OUT / "ablation_heads.png", 6.9, 1.4, width=6.0)
bullets(s, [
    ("Same backbone, same data, same schedule, 10 epochs; only the head differs.", 0),
    ("Band head (width-pooled, z-preserving): 4.96 mm val MAE.", 0),
    ("Standard global-average-pool + FC(15): 5.58 mm val MAE.", 0),
    ("The band head converges faster and ends ~11% better — consistent with the intuition that each output should read features from its own z-band.", 0),
    ("Both heads beat the 25.6 mm baseline by ~5× already after one epoch.", 0),
], left=0.6, top=1.5, width=6.1, size=16)

# ---------------- 9. Results ----------------
s = add_slide("Test-set results", "1000 held-out boards, evaluated once with the selected checkpoint")
rows = [
    ("Mean absolute error (the exercise metric)", "3.82 mm"),
    ("Median absolute error", "1.78 mm"),
    ("90th percentile absolute error", "7.67 mm"),
    ("Per-board MAE, median / worst", "1.83 / 116 mm"),
    ("Boards with pith fully inside the face (948)", "2.90 mm"),
    ("Boards with pith partly outside (52)", "20.6 mm"),
    ("Naive train-mean baseline", "25.6 mm"),
]
tbl = s.shapes.add_table(len(rows) + 1, 2, Inches(0.7), Inches(1.5), Inches(7.6), Inches(4.4)).table
tbl.columns[0].width = Inches(5.6)
tbl.columns[1].width = Inches(2.0)
hdr = ["Statistic", "Value"]
for j, t in enumerate(hdr):
    c = tbl.cell(0, j); c.text = t
    c.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
    c.text_frame.paragraphs[0].runs[0].font.bold = True
for i, (k, v) in enumerate(rows, 1):
    for j, t in enumerate([k, v]):
        c = tbl.cell(i, j); c.text = t
        run = c.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(14)
        if i == 1:
            run.font.bold = True
bullets(s, [
    ("MAE is uniform across the 15 z positions (3.6–4.2 mm; slightly higher at the board ends, where one-sided context is available).", 0),
    ("The error distribution is heavy-tailed: half of all boards are below 1.8 mm; a handful of far-outside-pith boards dominate the mean.", 0),
], left=8.5, top=1.7, width=4.4, size=14)

# ---------------- 10. Error analysis ----------------
s = add_slide("Where does the error come from?")
picture(s, OUT / "error_analysis_band_final.png", 0.35, 1.4, width=12.6)
bullets(s, [
    ("Left: most boards are almost perfectly predicted; a thin tail of hard boards remains.", 0),
    ("Right: error grows sharply once the true pith is beyond the board edges (dotted lines) — the model must extrapolate from nearly straight rings, and systematically underestimates large distances.", 0),
], top=5.1, size=15)

# ---------------- 11. Visual checks (typical) ----------------
s = add_slide("Visual check — randomly chosen test boards", "red = ground truth, blue = prediction")
picture(s, OUT / "visual_check_random_band_final.png", 0.35, 1.5, width=12.6)
bullets(s, [
    ("Predictions track ring symmetry closely, unaffected by knots and colour patches.", 0),
], top=5.3, size=15)

# ---------------- 12. Visual checks (outside + extremes) ----------------
s = add_slide("Visual check — pith outside the face, and failure cases")
picture(s, OUT / "visual_check_outside_band_final.png", 1.55, 1.3, width=10.2)
picture(s, OUT / "visual_check_extremes_band_final.png", 1.55, 4.15, width=10.2)
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run()
r.text = ("Top: pith just outside the face is handled well. Bottom right: the two worst test boards — "
          "pith 180–200 mm from the left edge; the model sees only near-vertical rings and underestimates the distance.")
r.font.size = Pt(13); r.font.color.rgb = GRAY

# ---------------- 13. Limitations ----------------
s = add_slide("Limitations")
bullets(s, [
    ("Far-outside piths are the dominant failure mode (52/1000 boards contribute a large share of the total error).", 0),
    ("Cause: weak visual signal (near-straight rings) + very few training examples that far out — an extrapolation problem in both image space and label space.", 1),
    ("Possible remedies: oversample/reweight these boards, a log-distance target parameterisation, or a heatmap head over an extended x-range.", 1),
    ("Synthetic data only: real board images add moisture stains, saw marks, lighting variation and camera geometry — domain gap not measured here.", 0),
    ("The 15 outputs are predicted without an explicit smoothness prior; the data show the true path wanders only ~10 mm, which could be exploited (e.g. low-order polynomial output or a smoothness penalty).", 0),
    ("Single-face x-only task: the full problem (x and y from four faces) needs multi-view fusion — the per-band head extends naturally to that setting.", 0),
    ("No uncertainty estimate: for grading decisions a confidence measure (e.g. quantile regression) would matter.", 0),
], top=1.35, size=16)

# ---------------- 14. Summary ----------------
s = add_slide("Summary — what I learned")
bullets(s, [
    ("A ~1 M-parameter CNN trained from scratch in ~1 hour on a laptop CPU reaches 3.82 mm test MAE — ~7× better than the naive baseline, ~2.6% of the board width.", 0),
    ("Understanding the data before modelling paid off:", 0),
    ("the smooth pith path and the 4.6% outside-the-board sub-population shaped the loss, the no-clamping output and the evaluation breakdown;", 1),
    ("physical symmetries gave exact, free augmentations.", 1),
    ("Encoding output structure in the architecture (width-pooled, z-preserving head) beat a generic regression head by ~11% at zero parameter cost — small inductive biases matter at this data scale.", 0),
    ("Honest evaluation = aggregate metric + stratified breakdown + visual checks; the mean alone hides that errors are concentrated in one rare, hard sub-population.", 0),
    ("Next steps: extended-range heatmap head for far piths, uncertainty estimation, and scaling to the full four-face x/y problem.", 0),
], top=1.35, size=17)

prs.save(OUT / "pith_position_presentation.pptx")
print("saved", OUT / "pith_position_presentation.pptx")
